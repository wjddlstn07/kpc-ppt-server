"""
server.py - KPC PPT 생성 Flask 서버
"""
import json
import os
import tempfile

from flask import Flask, jsonify, request, send_file
from ppt_builder import analyze_template, build_ppt
import anthropic

app = Flask(__name__)


# ──────────────────────────────────────────────
# 헬퍼: 슬라이드별 real_body 개수 추출
# ──────────────────────────────────────────────

def _extract_slide_structure(pptx_path: str) -> list:
    """
    템플릿 각 슬라이드의 real_body_candidates 개수를 반환.
    ppt_builder의 Shape 분류 로직과 동일한 기준 적용.
    """
    from pptx import Presentation
    from ppt_builder import (
        _find_title_shape,
        _find_pagenum_shape,
        _find_summary_shape,
        _has_real_text,
    )

    prs = Presentation(pptx_path)
    structure = []

    for i, slide in enumerate(prs.slides):
        text_shapes = [s for s in slide.shapes if s.has_text_frame]

        title_sh   = _find_title_shape(text_shapes)
        pagenum_sh = _find_pagenum_shape(text_shapes)
        remaining  = [s for s in text_shapes if s is not title_sh and s is not pagenum_sh]
        summary_sh = _find_summary_shape(remaining)

        body_candidates      = [s for s in remaining if s is not summary_sh]
        real_body_candidates = [s for s in body_candidates if _has_real_text(s)]

        structure.append({
            "slide_index"     : i,
            "real_body_count" : len(real_body_candidates),
            "real_body_names" : [s.name for s in real_body_candidates],
            "has_table"       : any(s.has_table for s in slide.shapes),
            "title_preview"   : title_sh.text.strip()[:30] if title_sh else "",
        })

    return structure


# ──────────────────────────────────────────────
# 라우트
# ──────────────────────────────────────────────

@app.route("/", methods=["GET"])
def health():
    return jsonify({"status": "ok", "message": "PPT 생성 서버 가동 중"})


@app.route("/analyze", methods=["POST"])
def analyze():
    """
    템플릿 .pptx 분석 → 슬라이드 구조 반환
    multipart/form-data: template (file)
    """
    if "template" not in request.files:
        return jsonify({"error": "template 파일이 필요합니다."}), 400

    template_file = request.files["template"]

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        template_path = tmp.name
        template_file.save(template_path)

    try:
        basic_info      = analyze_template(template_path)
        slide_structure = _extract_slide_structure(template_path)

        return jsonify({
            "basic_info"     : basic_info,
            "slide_structure": slide_structure,
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        os.unlink(template_path)


@app.route("/generate-ppt", methods=["POST"])
def generate_ppt():
    """
    slides_json 직접 입력 → PPTX 생성
    multipart/form-data: template (file), slides_json (text), title (text, optional)
    """
    if "template" not in request.files:
        return jsonify({"error": "template 파일이 필요합니다."}), 400
    if "slides_json" not in request.form:
        return jsonify({"error": "slides_json이 필요합니다."}), 400

    template_file = request.files["template"]
    slides_json   = request.form["slides_json"]
    title         = request.form.get("title", "output")

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        template_path = tmp.name
        template_file.save(template_path)

    output_path = tempfile.mktemp(suffix=".pptx")

    try:
        slides_data = json.loads(slides_json)
        build_ppt(template_path, slides_data, output_path)
        return send_file(
            output_path,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=f"{title}.pptx",
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        os.unlink(template_path)
        if os.path.exists(output_path):
            os.unlink(output_path)


@app.route("/generate-ppt-auto", methods=["POST"])
def generate_ppt_auto():
    """
    템플릿 구조 분석 → Claude JSON 생성 → PPTX 생성 (핵심 엔드포인트)
    multipart/form-data:
      - template     : .pptx 파일
      - content      : 프로젝트 내용 줄글
      - title        : 파일명 (선택)
    """
    if "template" not in request.files:
        return jsonify({"error": "template 파일이 필요합니다."}), 400
    if "content" not in request.form:
        return jsonify({"error": "content가 필요합니다."}), 400

    template_file = request.files["template"]
    content       = request.form["content"]
    title         = request.form.get("title", "output")

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        template_path = tmp.name
        template_file.save(template_path)

    output_path = tempfile.mktemp(suffix=".pptx")

    try:
        # ── 1. 템플릿 구조 분석
        slide_structure = _extract_slide_structure(template_path)
        print(f"[generate-ppt-auto] 템플릿 구조: {slide_structure}")

        # ── 2. 구조 기반 Claude 프롬프트 생성
        structure_desc = "\n".join([
            f"- 슬라이드 {s['slide_index'] + 1}: "
            f"본문 슬롯 {s['real_body_count']}개 "
            f"(제목 미리보기: '{s['title_preview']}')"
            for s in slide_structure
        ])

        bullet_rules = "\n".join([
            f"- 슬라이드 {s['slide_index'] + 1}: bullets 반드시 {s['real_body_count']}개"
            if s['real_body_count'] > 0
            else f"- 슬라이드 {s['slide_index'] + 1}: bullets 빈 배열 []"
            for s in slide_structure
        ])

        prompt = f"""아래 프로젝트 내용을 분석해서 PPT 슬라이드 구조를 JSON으로 만들어줘.
반드시 순수 JSON만 출력해. ```json 코드블록, 마크다운, 설명 텍스트 절대 포함하지 마. 첫 글자가 {{ 이고 마지막 글자가 }} 여야 해.

[템플릿 구조 - 반드시 준수]
이 템플릿은 {len(slide_structure)}개의 슬라이드로 구성됩니다.
{structure_desc}

[bullets 개수 규칙 - 절대 준수]
{bullet_rules}

[프로젝트 내용]
{content}

[출력 형식]
{{
  "slides": [
    {{
      "layout": "card",
      "title": "01 슬라이드 제목",
      "summary": "한 줄 요약",
      "bullets": ["항목1", "항목2", "항목3", "항목4"],
      "pageNum": 1
    }}
  ]
}}

[레이아웃 규칙]
- 비교·나열 → card
- 단계·순서 → process
- 데이터·실적 → table
- 범위·구조 → category
- 일정 → gantt (bullets 형식: "기간 | 단계 | 내용")
"""

        # ── 3. Claude API 호출
        client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))
        message = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=4096,
            messages=[{"role": "user", "content": prompt}],
        )

        raw_text = message.content[0].text.strip()
        print(f"[generate-ppt-auto] Claude 응답 미리보기: {raw_text[:200]}")

        # ── 4. JSON 파싱
        slides_data = json.loads(raw_text)

        # ── 5. PPTX 생성
        build_ppt(template_path, slides_data, output_path)

        return send_file(
            output_path,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=f"{title}.pptx",
        )
    except json.JSONDecodeError as e:
        return jsonify({"error": f"Claude JSON 파싱 실패: {str(e)}", "raw": raw_text}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        os.unlink(template_path)
        if os.path.exists(output_path):
            os.unlink(output_path)


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port, debug=False)
