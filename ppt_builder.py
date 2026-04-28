"""
ppt_builder.py - KPC 템플릿 기반 PPT 생성 모듈
"""
import copy

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn


# ──────────────────────────────────────────────
# 1. 템플릿 분석
# ──────────────────────────────────────────────

def analyze_template(pptx_path: str) -> dict:
    prs = Presentation(pptx_path)
    return {
        "slide_width_inches": round(prs.slide_width.inches, 2),
        "slide_height_inches": round(prs.slide_height.inches, 2),
        "layouts": _get_layouts(prs),
        "fonts": _get_fonts(prs),
        "existing_slides": _get_existing_slides(prs),
    }


def _get_layouts(prs: Presentation) -> list:
    layouts = []
    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = [
            {
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type).split(".")[-1],
                "name": ph.name,
            }
            for ph in layout.placeholders
        ]
        layouts.append({"index": idx, "name": layout.name, "placeholders": placeholders})
    return layouts


def _get_fonts(prs: Presentation) -> list:
    seen = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        seen.add(run.font.name)
    return sorted(seen)


def _get_existing_slides(prs: Presentation) -> list:
    EMU = 914400
    result = []
    for i, slide in enumerate(prs.slides):
        shapes_info = [
            {
                "name": s.name,
                "text_preview": s.text[:80].replace("\n", " "),
                "left_in": round(s.left / EMU, 2),
                "top_in": round(s.top / EMU, 2),
                "width_in": round(s.width / EMU, 2),
                "height_in": round(s.height / EMU, 2),
            }
            for s in slide.shapes
            if s.has_text_frame
        ]
        result.append({
            "index": i,
            "layout_name": slide.slide_layout.name,
            "shapes": shapes_info,
        })
    return result


# ──────────────────────────────────────────────
# 2. PPT 생성
# ──────────────────────────────────────────────

def build_ppt(template_path: str, slides_json: dict, output_path: str) -> str:
    prs = Presentation(template_path)
    slides_data = slides_json.get("slides", [])

    if not slides_data:
        raise ValueError("slides_json에 슬라이드 데이터가 없습니다.")

    template_count = len(prs.slides)
    if template_count == 0:
        raise ValueError("템플릿에 슬라이드가 없습니다.")

    print(f"[build_ppt] 템플릿 슬라이드 수: {template_count}, 생성할 슬라이드 수: {len(slides_data)}")

    tmpl_sp_trees = [copy.deepcopy(slide.shapes._spTree) for slide in prs.slides]

    _remove_all_slides(prs)

    for slide_idx, slide_data in enumerate(slides_data):
        tmpl_idx = min(slide_idx, template_count - 1)
        print(f"\n[슬라이드 {slide_idx + 1}] 템플릿 슬라이드 {tmpl_idx} 복사")
        slide = _add_slide_copy(prs, tmpl_sp_trees[tmpl_idx])
        _fill_slide(slide, slide_data)

    prs.save(output_path)
    return output_path


def _remove_all_slides(prs: Presentation) -> None:
    R_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    sld_id_lst = prs.slides._sldIdLst

    print(f"[_remove_all_slides] 제거 전 슬라이드 수: {len(prs.slides)}")

    for sld_id_elem in list(sld_id_lst):
        rId = sld_id_elem.get(R_ID)
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sld_id_lst.remove(sld_id_elem)

    print(f"[_remove_all_slides] 제거 후 슬라이드 수: {len(prs.slides)}")


def _add_slide_copy(prs: Presentation, tmpl_sp_tree) -> object:
    layout = _pick_layout(prs)
    new_slide = prs.slides.add_slide(layout)
    new_sp_tree = new_slide.shapes._spTree

    for child in list(new_sp_tree):
        new_sp_tree.remove(child)

    for child in tmpl_sp_tree:
        new_sp_tree.append(copy.deepcopy(child))

    return new_slide


def _pick_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if layout.placeholders:
            return layout
    return prs.slide_layouts[0]


# ──────────────────────────────────────────────
# 3. 슬라이드 내용 채우기
# ──────────────────────────────────────────────

def _has_real_text(shape) -> bool:
    """Shape에 실제 의미있는 텍스트가 있는지 확인 (공백·줄바꿈 제외, 최소 2자 이상)."""
    return len(shape.text.strip()) >= 2


def _fill_slide(slide, slide_data: dict) -> None:
    title_text   = slide_data.get("title", "")
    summary_text = slide_data.get("summary", "")
    bullets      = slide_data.get("bullets", [])
    page_num     = slide_data.get("pageNum")
    notes_text   = slide_data.get("notes", "")

    text_shapes = [s for s in slide.shapes if s.has_text_frame]

    title_sh   = _find_title_shape(text_shapes)
    pagenum_sh = _find_pagenum_shape(text_shapes)
    remaining  = [s for s in text_shapes if s is not title_sh and s is not pagenum_sh]
    summary_sh = _find_summary_shape(remaining)

    all_body_candidates = [s for s in remaining if s is not summary_sh]

    # ── 핵심 필터: 템플릿에 실제 텍스트가 있는 Shape만 콘텐츠 슬롯으로 인식
    real_body_candidates  = [s for s in all_body_candidates if _has_real_text(s)]
    empty_body_candidates = [s for s in all_body_candidates if not _has_real_text(s)]

    print(f"  title_shape         : {title_sh.name if title_sh else None}")
    print(f"  summary_shape       : {summary_sh.name if summary_sh else None}")
    print(f"  real_body_candidates: {[s.name for s in real_body_candidates]} ({len(real_body_candidates)}개)")
    print(f"  empty_body(장식용)  : {len(empty_body_candidates)}개 → 자동 clear")
    print(f"  pagenum_shape       : {pagenum_sh.name if pagenum_sh else None}")

    # ── 제목
    if title_sh and title_text:
        _replace_text(title_sh, title_text)
        print(f"  → title 채움: '{title_text[:40]}'")

    # ── 요약
    if summary_sh and summary_text:
        _replace_text(summary_sh, summary_text)
        print(f"  → summary 채움: '{summary_text[:40]}'")

    # ── 장식용 빈 Shape은 무조건 clear
    for s in empty_body_candidates:
        _clear_shape(s)

    # ── 본문: real_body_candidates 기준으로 매핑
    if bullets and real_body_candidates:
        n_bullets = len(bullets)
        n_shapes  = len(real_body_candidates)

        # 위→아래, 왼→오른 순 정렬
        sorted_bodies = sorted(
            real_body_candidates,
            key=lambda s: (round(s.top / 100000), s.left)
        )

        if n_bullets == n_shapes:
            print(f"  → 카드형 1:1 매핑 ({n_bullets}개)")
            for shape, bullet in zip(sorted_bodies, bullets):
                _replace_text(shape, bullet)
                print(f"     '{shape.name}' ← '{bullet[:30]}'")

        elif n_bullets < n_shapes:
            print(f"  → bullets({n_bullets}) < shapes({n_shapes}): 앞부터 채우고 나머지 clear")
            for i, shape in enumerate(sorted_bodies):
                if i < n_bullets:
                    _replace_text(shape, bullets[i])
                else:
                    _clear_shape(shape)

        else:
            body_sh = _find_body_shape(real_body_candidates)
            print(f"  → bullets({n_bullets}) > shapes({n_shapes}): 단일 body에 삽입")
            if body_sh:
                _replace_bullets(body_sh, bullets)
            for s in real_body_candidates:
                if s is not body_sh:
                    _clear_shape(s)

    elif not bullets and real_body_candidates:
        print(f"  → bullets 없음: real body shapes {len(real_body_candidates)}개 clear")
        for s in real_body_candidates:
            _clear_shape(s)

    # ── 페이지 번호
    if pagenum_sh and page_num is not None:
        _replace_text(pagenum_sh, str(page_num))
        print(f"  → pageNum 채움: {page_num}")

    if notes_text:
        try:
            slide.notes_slide.notes_text_frame.text = notes_text
        except Exception:
            pass


def _clear_shape(shape) -> None:
    """Shape 텍스트를 완전히 비웁니다 (첫 번째 단락 스타일은 유지)."""
    txBody = shape.text_frame._txBody
    paras  = txBody.findall(qn("a:p"))
    if not paras:
        return
    for p in paras[1:]:
        txBody.remove(p)
    for r in paras[0].findall(qn("a:r")):
        paras[0].remove(r)


# ──────────────────────────────────────────────
# 4. Shape 분류 헬퍼
# ──────────────────────────────────────────────

def _find_title_shape(text_shapes: list):
    if not text_shapes:
        return None

    for s in text_shapes:
        if any(kw in s.name.lower() for kw in ("title", "제목", "heading")):
            return s

    def _max_font_pt(shape):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    return run.font.size.pt
        return 0

    by_font = sorted(text_shapes, key=_max_font_pt, reverse=True)
    if _max_font_pt(by_font[0]) > 0:
        return by_font[0]

    return min(text_shapes, key=lambda s: s.top)


def _find_pagenum_shape(text_shapes: list):
    candidates = [
        s for s in text_shapes
        if len(s.text.strip()) <= 6 and any(c.isdigit() for c in s.text)
    ]
    if not candidates:
        return None
    return max(candidates, key=lambda s: s.top)


def _find_summary_shape(remaining: list):
    if not remaining:
        return None

    for s in remaining:
        if any(kw in s.name.lower() for kw in ("summary", "subtitle", "부제", "요약")):
            return s

    single_para = [s for s in remaining if len(s.text_frame.paragraphs) <= 2]
    if single_para:
        return min(single_para, key=lambda s: s.height)

    return None


def _find_body_shape(candidates: list):
    if not candidates:
        return None
    return max(candidates, key=lambda s: s.width * s.height)


# ──────────────────────────────────────────────
# 5. 텍스트 교체 (스타일 보존)
# ──────────────────────────────────────────────

def _replace_text(shape, new_text: str) -> None:
    if not new_text:
        return

    txBody = shape.text_frame._txBody
    paras = txBody.findall(qn("a:p"))
    if not paras:
        return

    first_p = paras[0]
    style_r = _capture_style_run(first_p)

    for p in paras[1:]:
        txBody.remove(p)

    for r in first_p.findall(qn("a:r")):
        first_p.remove(r)

    _append_run(first_p, new_text, style_r)


def _replace_bullets(shape, bullets: list) -> None:
    if not bullets:
        return

    txBody = shape.text_frame._txBody
    paras = txBody.findall(qn("a:p"))
    if not paras:
        return

    style_p = copy.deepcopy(paras[0])
    style_r = _capture_style_run(paras[0])

    for p in paras:
        txBody.remove(p)

    for bullet_text in bullets:
        new_p = copy.deepcopy(style_p)
        for r in new_p.findall(qn("a:r")):
            new_p.remove(r)
        _append_run(new_p, bullet_text, style_r)
        txBody.append(new_p)


def _capture_style_run(para_elem):
    runs = para_elem.findall(qn("a:r"))
    return copy.deepcopy(runs[0]) if runs else None


def _append_run(para_elem, text: str, style_r=None) -> None:
    if style_r is not None:
        new_r = copy.deepcopy(style_r)
        t_elem = new_r.find(qn("a:t"))
        if t_elem is None:
            t_elem = etree.SubElement(new_r, qn("a:t"))
        t_elem.text = text
        para_elem.append(new_r)
    else:
        r = etree.SubElement(para_elem, qn("a:r"))
        t = etree.SubElement(r, qn("a:t"))
        t.text = text
